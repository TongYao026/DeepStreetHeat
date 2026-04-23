import pptx

def debug_text(file_path):
    prs = pptx.Presentation(file_path)
    
    for i, slide in enumerate(prs.slides):
        slide_num = i + 1
        if slide_num in [10, 12, 17]:
            print(f"--- Slide {slide_num} ---")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    print(f"Shape Text: {shape.text}")

if __name__ == '__main__':
    debug_text('DeepStreetHeat-Multi-modal-Urban-Heat-Island-Analysis.pptx')
